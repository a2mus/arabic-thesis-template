"""Extract text from all Data/ source files for research ingestion.

Supported formats: PDF, DOCX, DOC (legacy), PPTX.
Outputs plain text files to extracted_texts/ for AI agent consumption.
"""
import os
import sys
import subprocess
import io

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")
sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding="utf-8", errors="replace")

DATA_DIR = os.path.join(os.path.dirname(__file__), "Data")
OUTPUT_DIR = os.path.join(os.path.dirname(__file__), "extracted_texts")
os.makedirs(OUTPUT_DIR, exist_ok=True)


def extract_pdf(filepath: str, output_path: str) -> None:
    """Extract text from PDF using pdfplumber."""
    try:
        import pdfplumber
    except ImportError:
        subprocess.check_call([sys.executable, "-m", "pip", "install", "pdfplumber", "-q"])
        import pdfplumber

    text_parts = []
    with pdfplumber.open(filepath) as pdf:
        for i, page in enumerate(pdf.pages):
            page_text = page.extract_text()
            if page_text:
                text_parts.append(f"--- Page {i+1} ---\n{page_text}")

    with open(output_path, "w", encoding="utf-8") as f:
        f.write("\n\n".join(text_parts))
    print(f"  PDF extracted: {len(text_parts)} pages -> {output_path}")


def extract_docx(filepath: str, output_path: str) -> None:
    """Extract text from DOCX using python-docx."""
    try:
        from docx import Document
    except ImportError:
        subprocess.check_call([sys.executable, "-m", "pip", "install", "python-docx", "-q"])
        from docx import Document

    doc = Document(filepath)
    text_parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            text_parts.append(para.text)

    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells)
            if row_text.strip():
                text_parts.append(f"[TABLE] {row_text}")

    with open(output_path, "w", encoding="utf-8") as f:
        f.write("\n".join(text_parts))
    print(f"  DOCX extracted: {len(text_parts)} paragraphs -> {output_path}")


def extract_doc(filepath: str, output_path: str) -> None:
    """Extract text from legacy .doc using binary parsing."""
    import re
    try:
        with open(filepath, "rb") as f:
            raw = f.read()
        text = raw.decode("utf-8", errors="ignore")
        text = re.sub(r'[^\u0600-\u06FF\u0020-\u007E\n\r\t\u060C\u061B\u061F]', ' ', text)
        text = re.sub(r'\s+', ' ', text).strip()

        with open(output_path, "w", encoding="utf-8") as f:
            f.write(text)
        print(f"  DOC extracted (raw): -> {output_path}")
    except Exception as e:
        print(f"  DOC extraction failed: {e}")


def extract_pptx(filepath: str, output_path: str) -> None:
    """Extract text from PPTX using python-pptx."""
    try:
        from pptx import Presentation
    except ImportError:
        subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx", "-q"])
        from pptx import Presentation

    prs = Presentation(filepath)
    text_parts = []
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text)
        if slide_texts:
            text_parts.append(f"--- Slide {slide_num} ---\n" + "\n".join(slide_texts))

    with open(output_path, "w", encoding="utf-8") as f:
        f.write("\n\n".join(text_parts))
    print(f"  PPTX extracted: {len(text_parts)} slides -> {output_path}")


def main() -> None:
    """Process all files in Data/ directory."""
    files = os.listdir(DATA_DIR)
    print(f"Found {len(files)} files in Data/\n")

    for filename in sorted(files):
        filepath = os.path.join(DATA_DIR, filename)
        if not os.path.isfile(filepath):
            continue

        base_name = os.path.splitext(filename)[0]
        ext = os.path.splitext(filename)[1].lower()
        safe_name = base_name.replace(" ", "_")
        output_path = os.path.join(OUTPUT_DIR, f"{safe_name}.txt")

        print(f"Processing: {filename}")

        if ext == ".pdf":
            extract_pdf(filepath, output_path)
        elif ext == ".docx":
            extract_docx(filepath, output_path)
        elif ext == ".doc":
            extract_doc(filepath, output_path)
        elif ext == ".pptx":
            extract_pptx(filepath, output_path)
        else:
            print(f"  Skipping unsupported format: {ext}")

        print()

    print("=== Extraction complete ===")
    for txt_file in sorted(os.listdir(OUTPUT_DIR)):
        txt_path = os.path.join(OUTPUT_DIR, txt_file)
        size = os.path.getsize(txt_path)
        print(f"  {txt_file}: {size:,} bytes")


if __name__ == "__main__":
    main()
